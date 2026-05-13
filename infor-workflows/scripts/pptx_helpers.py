"""python-pptx helpers shared across INFOR skills.

These helpers encode the formatting rules that recur across every INFOR deck
work and that have a history of regressing when re-derived inline in a skill:

  - `set_text(shape, lines)` preserves run-level rPr (font/size/bold/italic/color)
    by mutating runs[0].text in place rather than wiping and re-adding runs.
  - `write_bulleted_shape(shape, items)` harvests bullet pPr templates from the
    template's seed paragraphs BEFORE wiping, so new bullets keep the square /
    dash glyphs that python-pptx would otherwise drop.
  - `set_cell_text(cell, text, size_pt, color_hex)` forces Palatino on table
    cells (PowerPoint's default fallback is Calibri, which has been observed
    to slip in when cells are rewritten).
  - `clone_slide(prs, source_slide)` duplicates a slide and rewires every
    relationship (images, charts, hyperlinks) so the copy renders correctly
    instead of showing red-X placeholders.

Skills can load these via:

    import sys, os
    sys.path.insert(0, os.environ.get("CLAUDE_PLUGIN_ROOT", "./infor-workflows") + "/scripts")
    from pptx_helpers import set_text, write_bulleted_shape, set_cell_text, clone_slide, find_shape

Tests live next to this file in test_pptx_helpers.py and build fresh in-memory
decks so they don't depend on the INFOR template files.
"""

from copy import deepcopy

from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt


# ─── Brand constants ─────────────────────────────────────────────────────────

PALATINO = "Palatino Linotype"
COLOR_UP = "00B050"    # green — positive delta / beat
COLOR_DOWN = "C00000"  # red   — negative delta / miss


# ─── Shape lookup ────────────────────────────────────────────────────────────

def find_shape(slide, name):
    """Return the first shape on the slide whose .name matches."""
    for s in slide.shapes:
        if s.name == name:
            return s
    raise KeyError(f"Shape {name!r} not found on slide")


def find_shape_in_group(group, name):
    """Return the first child of the group whose .name matches."""
    for s in group.shapes:
        if s.name == name:
            return s
    raise KeyError(f"Shape {name!r} not found in group {group.name!r}")


# ─── XML helpers (private) ───────────────────────────────────────────────────

def _pPr_of(paragraph):
    """Return the paragraph-level <a:pPr> element, or None if absent."""
    for child in paragraph._p:
        if child.tag.endswith("}pPr"):
            return child
    return None


def _first_run_rPr(paragraph):
    """Return a deepcopy of the first run's <a:rPr>, or None if absent."""
    for r in paragraph.runs:
        for child in r._r:
            if child.tag.endswith("}rPr"):
                return deepcopy(child)
        return None
    return None


def _replace_run_rPr(run, template_rPr):
    """Strip the run's existing rPr children and graft a fresh deepcopy in."""
    if template_rPr is None:
        return
    for child in [c for c in run._r if c.tag.endswith("}rPr")]:
        run._r.remove(child)
    run._r.insert(0, deepcopy(template_rPr))


# ─── set_text ────────────────────────────────────────────────────────────────

def set_text(shape, lines, size_pt=None, color_hex=None):
    """Replace shape text while preserving the template's run formatting.

    Strategy:
      - For existing paragraphs (i < len(tf.paragraphs)): mutate runs[0].text
        in place to keep its rPr (font/size/bold/italic/color). Remove later
        runs on the same paragraph.
      - For new paragraphs beyond the existing count: add_paragraph, copy
        pPr from paragraph 0, copy first run's rPr from paragraph 0 so the
        new run inherits the template's formatting.

    `size_pt` and `color_hex` are explicit overrides applied AFTER the
    template formatting is restored. Pass them only when you intentionally
    want to override — e.g., the delta boxes on the earnings deck need
    forced 10 pt + green/red. Title bars / quote boxes should NOT receive
    overrides (the template's bold/italic/color would be wiped).
    """
    tf = shape.text_frame
    template_pPr = _pPr_of(tf.paragraphs[0])
    template_rPr = _first_run_rPr(tf.paragraphs[0])

    for i, line in enumerate(lines):
        if i < len(tf.paragraphs):
            p = tf.paragraphs[i]
            for r in list(p.runs[1:]):
                r._r.getparent().remove(r._r)
            if p.runs:
                # In-place mutation preserves the run's rPr
                p.runs[0].text = line
                run = p.runs[0]
            else:
                run = p.add_run()
                run.text = line
                _replace_run_rPr(run, template_rPr)
        else:
            p = tf.add_paragraph()
            if template_pPr is not None:
                for child in list(p._p):
                    if child.tag.endswith("}pPr"):
                        p._p.remove(child)
                p._p.insert(0, deepcopy(template_pPr))
            run = p.add_run()
            run.text = line
            _replace_run_rPr(run, template_rPr)

        if size_pt is not None:
            run.font.name = PALATINO
            run.font.size = Pt(size_pt)
        if color_hex is not None:
            run.font.color.rgb = RGBColor.from_string(color_hex)

    while len(tf.paragraphs) > len(lines):
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)


# ─── set_cell_text ───────────────────────────────────────────────────────────

def set_cell_text(cell, text, size_pt=9, color_hex=None):
    """Overwrite a table cell as a single Palatino run at size_pt.

    Unlike `set_text`, this DOES set font.name + font.size explicitly because
    PowerPoint's table-cell default fallback is Calibri — inheriting from the
    template has been observed to slip back to Calibri across rewrites.
    """
    tf = cell.text_frame
    while len(tf.paragraphs) > 1:
        last = tf.paragraphs[-1]
        last._p.getparent().remove(last._p)
    p = tf.paragraphs[0]
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = text
    run.font.name = PALATINO
    run.font.size = Pt(size_pt)
    if color_hex is not None:
        run.font.color.rgb = RGBColor.from_string(color_hex)


# ─── write_bulleted_shape ────────────────────────────────────────────────────

def _harvest_bullet_templates(shape):
    """Capture pPr + rPr templates from the shape's seed paragraphs BEFORE wiping.

    Returns {level_index: (pPr_copy, rPr_copy)} keyed 0, 1, 2... by ascending
    marL (smallest indent = main bullet = level 0).
    """
    tf = shape.text_frame
    harvested = []
    for para in tf.paragraphs:
        pPr = _pPr_of(para)
        rPr = _first_run_rPr(para)
        if pPr is None:
            continue
        marL = int(pPr.get("marL") or "0")
        harvested.append((marL, deepcopy(pPr), deepcopy(rPr) if rPr is not None else None))
    harvested.sort(key=lambda t: t[0])
    return {i: (pPr, rPr) for i, (_, pPr, rPr) in enumerate(harvested)}


def write_bulleted_shape(shape, items):
    """Wipe the shape and rewrite bullets with pPr + rPr correctly preserved.

    `items` is a list of tuples:
        (text, level)                       — single-run bullet
        (prefix_bold, rest_regular, level)  — two-run bullet (bold prefix +
                                              regular tail, used for segment
                                              names like 'easyfinancial: ...')

    `level` 0 = main (square glyph, larger font); 1 = sub (dash, smaller).

    Harvests the seed pPr templates BEFORE wiping so the bullet characters
    survive. Sets font.name = Palatino Linotype and font.size explicitly on
    every run. After writing, asserts every paragraph has a buChar element —
    raises RuntimeError if any bullet is missing its glyph, so a broken deck
    fails at write-time instead of shipping silently.
    """
    tf = shape.text_frame
    templates = _harvest_bullet_templates(shape)  # must happen BEFORE we wipe
    if not templates:
        raise RuntimeError(
            f"Shape {shape.name!r} has no bullet templates to harvest; "
            f"write_bulleted_shape requires the template to ship at least one "
            f"seed paragraph with a pPr (square or dash bullet)."
        )

    def _size_for(level):
        _, rPr = templates.get(level, (None, None))
        if rPr is not None:
            sz = rPr.get("sz")
            if sz is not None:
                return Pt(int(sz) / 100)
        return Pt(10.5 if level == 0 else 10.0)

    # Wipe: leave one paragraph behind, clear its runs and pPr
    while len(tf.paragraphs) > 1:
        last = tf.paragraphs[-1]
        last._p.getparent().remove(last._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    for child in list(first._p):
        if child.tag.endswith("}pPr"):
            first._p.remove(child)

    for i, item in enumerate(items):
        if len(item) == 2:
            prefix, rest, level = "", item[0], item[1]
        elif len(item) == 3:
            prefix, rest, level = item
        else:
            raise ValueError(
                "items must be (text, level) or (prefix_bold, rest_regular, level)"
            )

        p = first if i == 0 else tf.add_paragraph()
        if i != 0:
            for child in list(p._p):
                if child.tag.endswith("}pPr"):
                    p._p.remove(child)

        tmpl_pPr, _ = templates.get(level, templates[0])
        p._p.insert(0, deepcopy(tmpl_pPr))

        size = _size_for(level)
        if prefix:
            r1 = p.add_run()
            r1.text = prefix
            r1.font.name = PALATINO
            r1.font.size = size
            r1.font.bold = True
            r2 = p.add_run()
            r2.text = rest
            r2.font.name = PALATINO
            r2.font.size = size
            r2.font.bold = False
        else:
            r = p.add_run()
            r.text = rest
            r.font.name = PALATINO
            r.font.size = size
            r.font.bold = False

    # Post-write: every paragraph must have a bullet character
    for i, para in enumerate(tf.paragraphs):
        has_bu = False
        for elem in para._p.iter():
            if elem.tag.endswith("}buChar") or elem.tag.endswith("}buAutoNum"):
                has_bu = True
                break
        if not has_bu:
            raise RuntimeError(
                f"Shape {shape.name!r} paragraph {i} has no bullet character — "
                f"pPr template was not propagated. Refusing to ship a broken deck."
            )


# ─── Number formatting ───────────────────────────────────────────────────────

def fmt_broker_value(kind, value):
    """Format a broker table value with $ prefix or % suffix by metric kind.

    kind: 'dollar' | 'per_share' | 'percent' | 'volume'
    value: float (raw number) or str (already formatted; returned as-is).

    Dollar values use 1 decimal; per-share 2 decimals; percent 1 decimal +
    '%' suffix. Negatives are wrapped in parentheses, matching INFOR's
    financial-table convention.
    """
    if isinstance(value, str):
        return value
    neg = value < 0
    a = abs(value)
    if kind == "dollar":
        body = f"${a:,.1f}"
    elif kind == "per_share":
        body = f"${a:,.2f}"
    elif kind == "percent":
        return f"({a:,.1f}%)" if neg else f"{a:,.1f}%"
    elif kind == "volume":
        body = f"{a:,.1f}"  # caller appends a unit suffix if needed
    else:
        raise ValueError(f"unknown kind {kind!r}; expected dollar/per_share/percent/volume")
    return f"({body})" if neg else body


# ─── clone_slide ─────────────────────────────────────────────────────────────

# Shape XML attributes that carry relationship IDs — must be remapped when
# shapes are copied across slides, or pictures/charts/hyperlinks dangle.
_RID_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"))


def clone_slide(prs, source_slide):
    """Duplicate a slide, preserving shapes, layout, AND all relationships.

    Copying only the shape XML is not enough — shapes that reference
    relationships (pictures, charts, hyperlinks) carry r:embed / r:link / r:id
    attributes whose rIds point into the source slide's rels file. Without
    copying the rels and remapping rIds, those references dangle in the new
    slide and PowerPoint shows a red X with "The picture can't be displayed."

    The function:
      1. Adds a new slide using the source's layout (placeholders are wiped
         since we'll copy the source's own shapes in).
      2. Copies every non-notes relationship from source to new slide and
         records the old rId -> new rId mapping.
      3. Deep-copies each source shape's XML, walks every element, and
         rewrites r:embed / r:link / r:id attributes to the new rIds.

    Use this when building a deck from the INFOR Deck Template — clone the
    sample slide whose layout matches your content, then edit the clone.
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # Remove placeholders the layout auto-added — we want only the source's shapes
    for shp in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shp._element)

    # Copy every non-notes relationship and build an old -> new rId map
    rid_map = {}
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rid = new_slide.part.relate_to(
                rel.target_ref, rel.reltype, is_external=True
            )
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel.rId] = new_rid

    # Deep-copy each shape's XML and rewrite rId attributes to the new rIds
    for shp in source_slide.shapes:
        new_el = deepcopy(shp._element)
        for el in new_el.iter():
            for attr in _RID_ATTRS:
                if attr in el.attrib and el.attrib[attr] in rid_map:
                    el.attrib[attr] = rid_map[el.attrib[attr]]
        new_slide.shapes._spTree.append(new_el)

    return new_slide


def delete_slide(prs, index):
    """Remove a slide from a presentation by zero-based index.

    Used together with `clone_slide`: clone the sample slides you want, then
    delete the originals so only the analyst-edited clones remain.
    """
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])
