"""Tests for pptx_helpers.

Builds tiny in-memory decks with python-pptx — no INFOR template fixtures
required. Run with `python -m unittest infor-workflows/scripts/test_pptx_helpers.py`
or `python infor-workflows/scripts/test_pptx_helpers.py`.
"""

import os
import sys
import unittest
from copy import deepcopy

# Make pptx_helpers importable when run from the repo root or any cwd
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from pptx_helpers import (
    COLOR_DOWN,
    COLOR_UP,
    PALATINO,
    clone_slide,
    delete_slide,
    find_shape,
    find_shape_in_group,
    fmt_broker_value,
    set_cell_text,
    set_text,
    write_bulleted_shape,
)


# ─── Fixture builders ────────────────────────────────────────────────────────

def _make_shape_with_bold_italic_run(slide, name, text):
    """Add a textbox whose paragraph 0 has a bold-italic Palatino run — mimics
    the template's quote / title shapes that set_text must preserve."""
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.name = name
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = PALATINO
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.italic = True
    return tb


def _make_bulleted_shape(slide, name, seed_paragraphs):
    """Add a textbox with seed paragraphs that include explicit pPr + buChar,
    mimicking the template's bulleted shapes.

    seed_paragraphs: list of (text, marL_emu, font_size_pt, bullet_char)
    e.g. [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")]
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(4))
    tb.name = name
    tf = tb.text_frame

    # First paragraph already exists; reuse it for the first seed
    for i, (text, marL, size_pt, bu) in enumerate(seed_paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        # Wipe runs + pPr the python-pptx default left
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        for child in list(p._p):
            if child.tag.endswith("}pPr"):
                p._p.remove(child)

        # Build a <a:pPr marL="..." indent="..."><a:buChar char="..."/></a:pPr>
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pPr = etree.SubElement(p._p, f"{{{a_ns}}}pPr")
        pPr.set("marL", str(marL))
        pPr.set("indent", f"-{marL}")
        buChar = etree.SubElement(pPr, f"{{{a_ns}}}buChar")
        buChar.set("char", bu)
        # pPr must come BEFORE any run children — move it to the front
        p._p.remove(pPr)
        p._p.insert(0, pPr)

        run = p.add_run()
        run.text = text
        run.font.name = PALATINO
        run.font.size = Pt(size_pt)
    return tb


# ─── set_text tests ──────────────────────────────────────────────────────────

class TestSetText(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

    def test_preserves_bold_italic_when_no_overrides(self):
        """The whole point of set_text — mutating runs[0].text must keep the rPr.
        This is what would have shipped Calibri-default text if we'd wiped runs."""
        shape = _make_shape_with_bold_italic_run(self.slide, "Quote", "[original]")
        set_text(shape, ["the new quote"])
        run = shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.text, "the new quote")
        self.assertEqual(run.font.name, PALATINO)
        self.assertEqual(run.font.size, Pt(11))
        self.assertTrue(run.font.bold)
        self.assertTrue(run.font.italic)

    def test_size_override_applies(self):
        shape = _make_shape_with_bold_italic_run(self.slide, "Delta", "[x]")
        set_text(shape, ["+$4.2MM"], size_pt=10)
        run = shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.text, "+$4.2MM")
        self.assertEqual(run.font.size, Pt(10))

    def test_color_override_applies(self):
        shape = _make_shape_with_bold_italic_run(self.slide, "Delta", "[x]")
        set_text(shape, ["+$4.2MM"], size_pt=10, color_hex=COLOR_UP)
        run = shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(str(run.font.color.rgb), COLOR_UP)

    def test_multi_line_creates_new_paragraphs(self):
        shape = _make_shape_with_bold_italic_run(self.slide, "Footnote", "[seed]")
        set_text(shape, ["line 1", "line 2", "line 3"])
        tf = shape.text_frame
        self.assertEqual(len(tf.paragraphs), 3)
        self.assertEqual(tf.paragraphs[0].runs[0].text, "line 1")
        self.assertEqual(tf.paragraphs[1].runs[0].text, "line 2")
        self.assertEqual(tf.paragraphs[2].runs[0].text, "line 3")
        # New paragraph must inherit template formatting (italic + bold)
        new_run = tf.paragraphs[1].runs[0]
        self.assertEqual(new_run.font.name, PALATINO)
        self.assertTrue(new_run.font.italic)

    def test_shrinks_when_fewer_lines_than_paragraphs(self):
        shape = _make_shape_with_bold_italic_run(self.slide, "X", "[seed]")
        set_text(shape, ["a", "b", "c"])
        self.assertEqual(len(shape.text_frame.paragraphs), 3)
        set_text(shape, ["only one"])
        self.assertEqual(len(shape.text_frame.paragraphs), 1)
        self.assertEqual(shape.text_frame.paragraphs[0].runs[0].text, "only one")


# ─── set_cell_text tests ─────────────────────────────────────────────────────

class TestSetCellText(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.table = self.slide.shapes.add_table(
            rows=2, cols=2, left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(1)
        ).table

    def test_forces_palatino_at_size(self):
        set_cell_text(self.table.cell(0, 0), "Revenue", size_pt=9)
        run = self.table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.text, "Revenue")
        self.assertEqual(run.font.name, PALATINO)
        self.assertEqual(run.font.size, Pt(9))

    def test_color_hex_applies_to_variance_cell(self):
        set_cell_text(self.table.cell(1, 1), "+$1.2", size_pt=9, color_hex=COLOR_UP)
        run = self.table.cell(1, 1).text_frame.paragraphs[0].runs[0]
        self.assertEqual(str(run.font.color.rgb), COLOR_UP)

    def test_overwrite_replaces_previous_content(self):
        set_cell_text(self.table.cell(0, 0), "Old")
        set_cell_text(self.table.cell(0, 0), "New")
        self.assertEqual(self.table.cell(0, 0).text_frame.text, "New")


# ─── write_bulleted_shape tests ──────────────────────────────────────────────

class TestWriteBulletedShape(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

    def test_main_level_bullets_get_square_glyph(self):
        shape = _make_bulleted_shape(
            self.slide,
            "TextBox 16",
            [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")],
        )
        write_bulleted_shape(shape, [("bullet 1", 0), ("bullet 2", 0)])

        paras = shape.text_frame.paragraphs
        self.assertEqual(len(paras), 2)
        for p in paras:
            buChars = [e for e in p._p.iter() if e.tag.endswith("}buChar")]
            self.assertEqual(len(buChars), 1, "every paragraph must have a buChar")
            self.assertEqual(buChars[0].get("char"), "■", "level-0 square glyph")

    def test_sub_level_bullets_get_dash_glyph(self):
        shape = _make_bulleted_shape(
            self.slide,
            "TextBox 16",
            [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")],
        )
        write_bulleted_shape(shape, [("main", 0), ("sub one", 1), ("sub two", 1)])

        paras = shape.text_frame.paragraphs
        glyphs = []
        for p in paras:
            buChar = next(e for e in p._p.iter() if e.tag.endswith("}buChar"))
            glyphs.append(buChar.get("char"))
        self.assertEqual(glyphs, ["■", "-", "-"])

    def test_bold_prefix_two_run_paragraph(self):
        shape = _make_bulleted_shape(
            self.slide,
            "TextBox 16",
            [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")],
        )
        write_bulleted_shape(
            shape,
            [("easyfinancial:", " direct-to-consumer unsecured loans", 1)],
        )
        p = shape.text_frame.paragraphs[0]
        self.assertEqual(len(p.runs), 2)
        self.assertEqual(p.runs[0].text, "easyfinancial:")
        self.assertTrue(p.runs[0].font.bold)
        self.assertEqual(p.runs[1].text, " direct-to-consumer unsecured loans")
        self.assertFalse(p.runs[1].font.bold)

    def test_explicit_palatino_on_every_run(self):
        shape = _make_bulleted_shape(
            self.slide,
            "TextBox 16",
            [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")],
        )
        write_bulleted_shape(shape, [("a", 0), ("b", 1), ("c:", " tail", 0)])
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                self.assertEqual(
                    r.font.name, PALATINO,
                    f"every run must have Palatino set explicitly (got {r.font.name!r})"
                )

    def test_raises_when_template_has_no_pPr(self):
        """Defensive: a shape with no seed bullets cannot be safely rewritten."""
        tb = self.slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.name = "Empty"
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "no bullet here"

        with self.assertRaises(RuntimeError):
            write_bulleted_shape(tb, [("x", 0)])

    def test_bad_item_shape_raises(self):
        shape = _make_bulleted_shape(
            self.slide, "TextBox 16",
            [("seed", 180975, 10.5, "■")],
        )
        with self.assertRaises(ValueError):
            write_bulleted_shape(shape, [("a", "b", "c", "d")])  # wrong arity


# ─── fmt_broker_value tests ──────────────────────────────────────────────────

class TestFmtBrokerValue(unittest.TestCase):
    def test_dollar_positive(self):
        self.assertEqual(fmt_broker_value("dollar", 406.3), "$406.3")

    def test_dollar_negative_uses_parens(self):
        self.assertEqual(fmt_broker_value("dollar", -121.1), "($121.1)")

    def test_dollar_thousands_separator(self):
        self.assertEqual(fmt_broker_value("dollar", 1234.5), "$1,234.5")

    def test_per_share_two_decimals(self):
        self.assertEqual(fmt_broker_value("per_share", 1.24), "$1.24")
        self.assertEqual(fmt_broker_value("per_share", -8.93), "($8.93)")

    def test_percent_with_suffix(self):
        self.assertEqual(fmt_broker_value("percent", 38.7), "38.7%")
        self.assertEqual(fmt_broker_value("percent", -2.3), "(2.3%)")

    def test_already_formatted_string_passes_through(self):
        self.assertEqual(fmt_broker_value("dollar", "$406.3"), "$406.3")
        self.assertEqual(fmt_broker_value("percent", "N/A"), "N/A")

    def test_unknown_kind_raises(self):
        with self.assertRaises(ValueError):
            fmt_broker_value("oranges", 5.0)


# ─── find_shape / find_shape_in_group tests ──────────────────────────────────

class TestFindShape(unittest.TestCase):
    def test_find_shape_by_name(self):
        prs = Presentation()
        # Layout 6 is "Blank" — no auto-added placeholders to collide with our textbox name
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tb.name = "Rectangle 1032"
        found = find_shape(slide, "Rectangle 1032")
        # python-pptx returns a fresh proxy each iteration — equality is the right check
        self.assertEqual(found.name, "Rectangle 1032")
        self.assertEqual(found.shape_id, tb.shape_id)

    def test_find_shape_missing_raises(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        with self.assertRaises(KeyError):
            find_shape(slide, "Nonexistent")


# ─── clone_slide / delete_slide tests ────────────────────────────────────────

class TestCloneSlide(unittest.TestCase):
    def _make_source_slide_with_shapes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        # Add two distinguishable shapes
        tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tb1.name = "Source TextBox A"
        tb1.text_frame.paragraphs[0].add_run().text = "alpha"
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        tb2.name = "Source TextBox B"
        tb2.text_frame.paragraphs[0].add_run().text = "beta"
        return prs, slide, [tb1.name, tb2.name]

    def test_clone_copies_shape_xml(self):
        prs, source, names = self._make_source_slide_with_shapes()
        new_slide = clone_slide(prs, source)
        new_names = [s.name for s in new_slide.shapes]
        for name in names:
            self.assertIn(name, new_names)

    def test_clone_preserves_text(self):
        prs, source, _ = self._make_source_slide_with_shapes()
        new_slide = clone_slide(prs, source)
        texts = [s.text_frame.text for s in new_slide.shapes if s.has_text_frame]
        self.assertIn("alpha", texts)
        self.assertIn("beta", texts)

    def test_clone_inherits_layout(self):
        prs, source, _ = self._make_source_slide_with_shapes()
        new_slide = clone_slide(prs, source)
        self.assertEqual(new_slide.slide_layout.name, source.slide_layout.name)

    def test_clone_picture_remaps_rid(self):
        """A picture has an r:embed rId — the new slide must own a fresh rId
        pointing at the same image so the picture renders instead of red-X."""
        import io
        # Tiny valid PNG (1x1 black pixel)
        png_bytes = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cc\xf8"
            b"\xcf\xc0\x00\x00\x00\x03\x00\x01\\\xcd\xff\x69\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        prs = Presentation()
        source = prs.slides.add_slide(prs.slide_layouts[6])
        source.shapes.add_picture(io.BytesIO(png_bytes), Inches(1), Inches(1), Inches(1), Inches(1))

        new_slide = clone_slide(prs, source)

        # Find the picture on the new slide and confirm its r:embed rId
        # resolves to an existing relationship (would raise KeyError if not)
        found_picture = False
        for shp in new_slide.shapes:
            for el in shp._element.iter():
                embed = el.get(qn("r:embed"))
                if embed:
                    found_picture = True
                    self.assertIn(embed, new_slide.part.rels,
                                  "cloned picture's r:embed must point at a new-slide rel")
        self.assertTrue(found_picture, "cloned slide must contain the picture")


class TestDeleteSlide(unittest.TestCase):
    def test_delete_removes_slide(self):
        prs = Presentation()
        s0 = prs.slides.add_slide(prs.slide_layouts[6])
        s0.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).name = "keep-me"
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).name = "delete-me"
        self.assertEqual(len(prs.slides), 2)
        delete_slide(prs, 1)
        self.assertEqual(len(prs.slides), 1)
        remaining = [s.name for s in prs.slides[0].shapes]
        self.assertIn("keep-me", remaining)


# ─── Constants ───────────────────────────────────────────────────────────────

class TestConstants(unittest.TestCase):
    def test_palatino_full_name(self):
        self.assertEqual(PALATINO, "Palatino Linotype")

    def test_brand_color_hex_uppercase_no_hash(self):
        self.assertEqual(COLOR_UP, "00B050")
        self.assertEqual(COLOR_DOWN, "C00000")


if __name__ == "__main__":
    unittest.main(verbosity=2)
